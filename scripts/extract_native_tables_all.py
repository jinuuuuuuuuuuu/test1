import io
import os
import re
import glob
import openpyxl

TRACKER = r"C:\Users\kevin\Downloads\docs 수정.xlsm"
DOCS_DIR = r"C:\Users\kevin\pension-agent\data\raw\docs\docs_renamed"
OUT = r"C:\Users\kevin\pension-agent\docs_check\_native_tables_extracted.txt"

wb = openpyxl.load_workbook(TRACKER, data_only=True)
ws = wb["파싱 결과"]

docnums = set()
for row in ws.iter_rows(min_row=6):
    if str(row[9].value).strip().upper() == "Y":
        m = re.match(r"doc(\d+)_chunk\d+", str(row[7].value))
        if m:
            docnums.add(int(m.group(1)))

docnums = sorted(docnums)
print("target docs:", docnums)

results = {}

for num in docnums:
    matches = glob.glob(os.path.join(DOCS_DIR, f"doc{num}.*"))
    matches = [m for m in matches if not m.endswith(".txt")]
    if not matches:
        results[num] = ("NOFILE", [])
        continue
    path = matches[0]
    ext = os.path.splitext(path)[1].lower()

    tables_out = []
    try:
        if ext == ".pdf":
            import fitz
            d = fitz.open(path)
            for pno in range(d.page_count):
                page = d[pno]
                tabs = page.find_tables()
                for ti, t in enumerate(tabs.tables):
                    grid = t.extract()
                    tables_out.append((f"page{pno+1}_table{ti+1}", grid))
            d.close()
        elif ext == ".pptx":
            from pptx import Presentation
            p = Presentation(path)
            for si, slide in enumerate(p.slides):
                ti = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        ti += 1
                        grid = []
                        for r in shape.table.rows:
                            grid.append([c.text for c in r.cells])
                        tables_out.append((f"slide{si+1}_table{ti}", grid))
        elif ext == ".docx":
            import docx
            d = docx.Document(path)
            for ti, table in enumerate(d.tables):
                grid = []
                for r in table.rows:
                    grid.append([c.text for c in r.cells])
                tables_out.append((f"table{ti+1}", grid))
        elif ext == ".xlsx":
            xwb = openpyxl.load_workbook(path, data_only=True)
            for sheetname in xwb.sheetnames:
                sh = xwb[sheetname]
                grid = []
                for r in sh.iter_rows(values_only=True):
                    if any(c is not None for c in r):
                        grid.append([str(c) if c is not None else "" for c in r])
                if grid:
                    tables_out.append((f"sheet[{sheetname}]", grid))
        else:
            results[num] = (f"UNSUPPORTED:{ext}", [])
            continue
    except Exception as e:
        results[num] = (f"ERROR:{e}", [])
        continue

    results[num] = (ext, tables_out)

with io.open(OUT, "w", encoding="utf-8") as f:
    for num in docnums:
        ext, tables_out = results[num]
        f.write(f"########## doc{num} ({ext}) — {len(tables_out)} table(s) found ##########\n")
        for label, grid in tables_out:
            f.write(f"--- {label}: {len(grid)} rows x {len(grid[0]) if grid else 0} cols ---\n")
            for r in grid:
                f.write(" | ".join((c or "").replace("\n", " ¶ ") for c in r) + "\n")
            f.write("\n")
        f.write("\n")

print("done ->", OUT)
