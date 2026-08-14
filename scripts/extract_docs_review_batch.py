"""검수전 372개 청크가 있는 40개 문서의 트래커 청크 데이터 + 원본 전체 텍스트를 뽑아
docs_check/{docN}.txt 로 저장한다. (pdf/docx/pptx/xlsx 전부 지원)"""

import os
import openpyxl
import fitz
import docx
import pptx

TRACKER_PATH = r"C:\Users\kevin\Downloads\docs 수정.xlsm"
DOCS_DIR = r"C:\Users\kevin\pension-agent\data\raw\docs\docs_renamed"
OUT_DIR = r"C:\Users\kevin\pension-agent\docs_check"

TARGET_DOCS = [
    6, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 22, 23, 25, 26, 29, 33, 34, 35, 36,
    38, 39, 40, 41, 42, 43, 44, 45, 46, 47, 48, 49, 50, 51, 52, 53, 55, 57, 58,
]

os.makedirs(OUT_DIR, exist_ok=True)

wb = openpyxl.load_workbook(TRACKER_PATH, keep_vba=True, data_only=True)
ws = wb["파싱 결과"]

header = ['파일번호', '파일제목', '파일형식', '카테고리', '저장방식', 'section', '원문위치',
          'chunk_id', 'chunk_text', '표포함', '표구조화텍스트', '특이사항', '검수상태']

rows_by_doc = {}
for row in ws.iter_rows(min_row=6, values_only=True):
    fnum = row[0]
    if fnum is None:
        continue
    try:
        num = int(float(fnum))
    except (TypeError, ValueError):
        continue
    if num in TARGET_DOCS:
        rows_by_doc.setdefault(num, []).append(row)


def extract_pdf(path):
    doc = fitz.open(path)
    parts = []
    for i in range(doc.page_count):
        parts.append(f"\n===== PAGE {i+1}/{doc.page_count} =====\n")
        parts.append(doc[i].get_text())
    doc.close()
    return "".join(parts)


def extract_docx(path):
    d = docx.Document(path)
    parts = []
    for i, para in enumerate(d.paragraphs):
        if para.text.strip():
            parts.append(para.text)
    for ti, table in enumerate(d.tables):
        parts.append(f"\n===== TABLE {ti+1} =====")
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def extract_pptx(path):
    p = pptx.Presentation(path)
    parts = []
    for i, slide in enumerate(p.slides):
        parts.append(f"\n===== SLIDE {i+1} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def extract_xlsx(path):
    wb2 = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb2.worksheets:
        parts.append(f"\n===== SHEET: {sheet.title} =====")
        for row in sheet.iter_rows(values_only=True):
            if any(v is not None and str(v).strip() for v in row):
                parts.append(" | ".join("" if v is None else str(v) for v in row))
    return "\n".join(parts)


EXTRACTORS = {".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx, ".xlsx": extract_xlsx}

for num in TARGET_DOCS:
    out_path = os.path.join(OUT_DIR, f"doc{num}.txt")
    src_path = None
    for ext in EXTRACTORS:
        candidate = os.path.join(DOCS_DIR, f"doc{num}{ext}")
        if os.path.exists(candidate):
            src_path = candidate
            break
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"########## TRACKER DATA: doc{num} ##########\n")
        for row in rows_by_doc.get(num, []):
            f.write(f"--- {row[7]} (status={row[12]}) ---\n")
            for i, h in enumerate(header):
                val = row[i]
                if val is not None and str(val).strip() != "":
                    f.write(f"  [{h}] {val}\n")
        f.write("\n########## SOURCE FULL TEXT ##########\n")
        if not src_path:
            f.write("*** FILE NOT FOUND ***\n")
            print(f"doc{num}: NOT FOUND")
            continue
        ext = os.path.splitext(src_path)[1]
        try:
            text = EXTRACTORS[ext](src_path)
        except Exception as e:
            text = f"*** EXTRACTION ERROR: {e} ***"
        f.write(text)
    print(f"doc{num} -> {out_path}")

print("Done:", len(TARGET_DOCS), "docs")
