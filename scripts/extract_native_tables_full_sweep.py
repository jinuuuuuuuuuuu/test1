import io
import os
import re
import glob
import openpyxl

TRACKER = r"C:\Users\kevin\Downloads\docs 수정.xlsm"
DOCS_DIR = r"C:\Users\kevin\pension-agent\data\raw\docs\docs_renamed"
OUT = r"C:\Users\kevin\pension-agent\docs_check\_native_tables_full_sweep.txt"
SUMMARY_OUT = r"C:\Users\kevin\pension-agent\docs_check\_native_tables_full_sweep_summary.txt"

wb = openpyxl.load_workbook(TRACKER, data_only=True)
ws = wb["파싱 결과"]

# doc번호 -> {있는 chunk_id 목록, Y로 표시된 chunk_id 목록}
doc_chunks = {}
for row in ws.iter_rows(min_row=6):
    cid = row[7].value
    if not cid:
        continue
    m = re.match(r"doc(\d+)_chunk(\d+)", str(cid))
    if not m:
        continue
    num = int(m.group(1))
    has_table = str(row[9].value).strip().upper() == "Y"
    doc_chunks.setdefault(num, {"all": [], "Y": []})
    doc_chunks[num]["all"].append(cid)
    if has_table:
        doc_chunks[num]["Y"].append(cid)

all_docnums = sorted(doc_chunks.keys())
print("total docs in tracker:", len(all_docnums))

results = {}

for num in all_docnums:
    matches = glob.glob(os.path.join(DOCS_DIR, f"doc{num}.*"))
    matches = [m for m in matches if not m.endswith(".txt")]
    if not matches:
        results[num] = ("NOFILE", [])
        continue
    path = matches[0]
    ext = os.path.splitext(path)[1].lower()

    tables_out = []
    try:
        if ext == ".pdf":
            import fitz
            d = fitz.open(path)
            for pno in range(d.page_count):
                page = d[pno]
                tabs = page.find_tables()
                for ti, t in enumerate(tabs.tables):
                    grid = t.extract()
                    # 실질적 표만(2행 이상 또는 2열 이상 & 내용 있음)
                    nonempty = sum(1 for r in grid for c in r if c and str(c).strip())
                    if nonempty >= 2:
                        tables_out.append((f"page{pno+1}_table{ti+1}", grid))
            d.close()
        elif ext == ".pptx":
            from pptx import Presentation
            p = Presentation(path)
            for si, slide in enumerate(p.slides):
                ti = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        ti += 1
                        grid = [[c.text for c in r.cells] for r in shape.table.rows]
                        tables_out.append((f"slide{si+1}_table{ti}", grid))
        elif ext == ".docx":
            import docx
            d = docx.Document(path)
            for ti, table in enumerate(d.tables):
                grid = [[c.text for c in r.cells] for r in table.rows]
                tables_out.append((f"table{ti+1}", grid))
        elif ext == ".xlsx":
            xwb = openpyxl.load_workbook(path, data_only=True)
            for sheetname in xwb.sheetnames:
                sh = xwb[sheetname]
                grid = []
                for r in sh.iter_rows(values_only=True):
                    if any(c is not None for c in r):
                        grid.append([str(c) if c is not None else "" for c in r])
                if grid:
                    tables_out.append((f"sheet[{sheetname}]", grid))
        else:
            results[num] = (f"UNSUPPORTED:{ext}", [])
            continue
    except Exception as e:
        results[num] = (f"ERROR:{e}", [])
        continue

    results[num] = (ext, tables_out)

# 상세 덤프
with io.open(OUT, "w", encoding="utf-8") as f:
    for num in all_docnums:
        ext, tables_out = results[num]
        y_count = len(doc_chunks[num]["Y"])
        f.write(f"########## doc{num} ({ext}) — native table {len(tables_out)}건 / 표포함=Y 청크 {y_count}건 ##########\n")
        for label, grid in tables_out:
            f.write(f"--- {label}: {len(grid)}행 x {len(grid[0]) if grid else 0}열 ---\n")
            for r in grid:
                f.write(" | ".join((c or "").replace("\n", " ¶ ") for c in r) + "\n")
            f.write("\n")
        f.write("\n")

# 불일치 요약: native 표 있는데 Y태그 없음 / Y태그 있는데 native 표 없음
with io.open(SUMMARY_OUT, "w", encoding="utf-8") as f:
    f.write(f"전체 문서: {len(all_docnums)}건\n\n")
    f.write("=== A. 원문에 표(native table)가 있는데 표포함=Y로 태그된 청크가 없는 문서 ===\n")
    for num in all_docnums:
        ext, tables_out = results[num]
        y_count = len(doc_chunks[num]["Y"])
        if tables_out and y_count == 0:
            f.write(f"doc{num} ({ext}): native 표 {len(tables_out)}건, Y태그 0건\n")

    f.write("\n=== B. 표포함=Y 태그는 있는데 native 표가 0건인 문서 (PDF 인코딩 문제 또는 이미지형일 수 있음, 기존 목록과 대조 필요) ===\n")
    for num in all_docnums:
        ext, tables_out = results[num]
        y_count = len(doc_chunks[num]["Y"])
        if y_count > 0 and not tables_out:
            f.write(f"doc{num} ({ext}): native 표 0건, Y태그 {y_count}건 -> {doc_chunks[num]['Y']}\n")

    f.write("\n=== C. 파일 없음/에러 ===\n")
    for num in all_docnums:
        ext, tables_out = results[num]
        if ext.startswith("NOFILE") or ext.startswith("ERROR") or ext.startswith("UNSUPPORTED"):
            f.write(f"doc{num}: {ext}\n")

print("done ->", OUT, SUMMARY_OUT)
